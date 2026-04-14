"""
Script to extract text content from PowerPoint (.pptx) files.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pathlib import Path
import sys


def extract_text_from_pptx(pptx_path: str) -> dict:
    """
    Extract all text content from a PowerPoint file.
    
    Args:
        pptx_path: Path to the .pptx file
        
    Returns:
        Dictionary with slide numbers as keys and list of text content as values
    """
    prs = Presentation(pptx_path)
    slides_content = {}
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            
            # Extract text from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        slide_texts.append(" | ".join(row_text))
        
        slides_content[slide_num] = slide_texts
    
    return slides_content


def extract_to_string(pptx_path: str) -> str:
    """
    Extract text from PowerPoint and return as formatted string.
    """
    content = extract_text_from_pptx(pptx_path)
    output_lines = []
    
    filename = Path(pptx_path).name
    output_lines.append(f"=== Text extracted from: {filename} ===\n")
    
    for slide_num, texts in content.items():
        output_lines.append(f"\n--- Slide {slide_num} ---")
        for text in texts:
            output_lines.append(text)
    
    return "\n".join(output_lines)


def save_extracted_text(pptx_path: str, output_path: str = None) -> str:
    """
    Extract text from PowerPoint and save to a text file.
    
    Args:
        pptx_path: Path to the .pptx file
        output_path: Optional output path. If None, creates .txt file with same name.
        
    Returns:
        Path to the output file
    """
    if output_path is None:
        output_path = Path(pptx_path).with_suffix('.txt')
    
    text_content = extract_to_string(pptx_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text_content)
    
    return str(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx_text.py <path_to_pptx> [output_path]")
        print("\nExample:")
        print("  python extract_pptx_text.py presentation.pptx")
        print("  python extract_pptx_text.py presentation.pptx output.txt")
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    if not Path(pptx_file).exists():
        print(f"Error: File not found: {pptx_file}")
        sys.exit(1)
    
    if output_file:
        saved_path = save_extracted_text(pptx_file, output_file)
        print(f"Text extracted and saved to: {saved_path}")
    else:
        # Print to console
        print(extract_to_string(pptx_file))
